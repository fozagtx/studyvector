#!/usr/bin/env python3
"""
StudyVector — Ingestion Pipeline
Scans materials/ for PDFs, PPTX, and DOC/DOCX files, extracts text, chunks it, embeds it, stores in VectorAI DB.

Usage:
    python ingest.py                        # ingest everything in ./materials
    python ingest.py --dir ~/Downloads/CS   # custom folder
    python ingest.py --reset                # wipe and re-ingest from scratch
"""

import argparse
import hashlib
import json
import os
import re
import sys
from pathlib import Path

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from tqdm import tqdm

try:
    from actian_vectorai import VectorAIClient, VectorParams, Distance, PointStruct
except ImportError:
    print("ERROR: actian-vectorai not installed.")
    print("  pip install actian_vectorai-0.1.0b2-py3-none-any.whl")
    sys.exit(1)

# ── Config ──────────────────────────────────────────────────────────────────
COLLECTION     = "studyvector"
EMBED_MODEL    = "all-MiniLM-L6-v2"   # 384-dim, fast, fully offline
VECTOR_DIM     = 384
CHUNK_SIZE     = 400   # characters
CHUNK_OVERLAP  = 80
SERVER         = "localhost:50051"
PROGRESS_FILE  = ".ingest_progress.json"  # tracks already-ingested files


# ── Text extraction ──────────────────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".doc", ".docx"}


def extract_file(path: Path) -> list[dict]:
    """Return list of {page, text} dicts from a supported file."""
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(path)
    elif ext == ".pptx":
        return _extract_pptx(path)
    elif ext in (".doc", ".docx"):
        return _extract_docx(path)
    return []


def _extract_pdf(path: Path) -> list[dict]:
    """Return list of {page, text} dicts from a PDF."""
    pages = []
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                text = text.strip()
                if text:
                    pages.append({"page": i, "text": text})
    except Exception as e:
        print(f"  WARNING: could not read {path.name}: {e}")
    return pages


def _extract_pptx(path: Path) -> list[dict]:
    """Return list of {page, text} dicts from a PPTX (one entry per slide)."""
    pages = []
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            texts.append(line)
            text = "\n".join(texts)
            if text:
                pages.append({"page": i, "text": text})
    except Exception as e:
        print(f"  WARNING: could not read {path.name}: {e}")
    return pages


def _extract_docx(path: Path) -> list[dict]:
    """Return list of {page, text} dicts from a DOCX (.doc not supported natively)."""
    pages = []
    try:
        doc = DocxDocument(path)
        full_text = "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
        if full_text:
            pages.append({"page": 1, "text": full_text})
    except Exception as e:
        print(f"  WARNING: could not read {path.name}: {e}")
    return pages


# ── Chunking ─────────────────────────────────────────────────────────────────

def chunk_pages(pages: list[dict], source: str, course: str) -> list[dict]:
    """Split page texts into overlapping chunks with metadata."""
    chunks = []
    for page_info in pages:
        text = page_info["text"]
        page_num = page_info["page"]
        start = 0
        while start < len(text):
            end = min(start + CHUNK_SIZE, len(text))
            # snap to word boundary
            if end < len(text):
                snap = text.rfind(" ", start, end)
                if snap > start:
                    end = snap
            chunk = text[start:end].strip()
            if len(chunk) > 40:  # skip tiny fragments
                chunks.append({
                    "text":    chunk,
                    "source":  source,   # filename
                    "course":  course,   # derived from folder name
                    "page":    page_num,
                })
            start = end - CHUNK_OVERLAP if end < len(text) else len(text)
    return chunks


def derive_course(path: Path, root: Path) -> str:
    """Use parent folder name as course label, falling back to 'General'."""
    rel = path.relative_to(root)
    parts = rel.parts
    if len(parts) > 1:
        return parts[0]
    return "General"


# ── Ingestion ────────────────────────────────────────────────────────────────

def load_progress() -> dict:
    if Path(PROGRESS_FILE).exists():
        with open(PROGRESS_FILE) as f:
            return json.load(f)
    return {}


def save_progress(progress: dict):
    with open(PROGRESS_FILE, "w") as f:
        json.dump(progress, f, indent=2)


def file_hash(path: Path) -> str:
    h = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def make_client():
    """Create a fresh VectorAI client (avoids gRPC keepalive timeouts)."""
    c = VectorAIClient(SERVER)
    c.connect()
    return c


def ingest(materials_dir: Path, reset: bool = False):
    files = sorted(
        f for f in materials_dir.rglob("*")
        if f.suffix.lower() in SUPPORTED_EXTENSIONS
    )
    if not files:
        print(f"No supported files found in {materials_dir}. Drop your materials there and re-run.")
        return

    print(f"Found {len(files)} file(s) in {materials_dir}")

    print(f"\nLoading embedding model ({EMBED_MODEL})...")
    model = SentenceTransformer(EMBED_MODEL)

    # Quick setup — use a short-lived connection for collection management
    print(f"Connecting to VectorAI DB at {SERVER}...")
    client = make_client()
    info = client.health_check()
    print(f"  Connected: {info['title']} v{info['version']}")

    exists = client.collections.exists(COLLECTION)
    if reset and exists:
        client.collections.delete(COLLECTION)
        print(f"  Wiped existing collection '{COLLECTION}'")
        exists = False
        save_progress({})

    if not exists:
        client.collections.create(
            COLLECTION,
            vectors_config=VectorParams(size=VECTOR_DIM, distance=Distance.Cosine),
        )
        print(f"  Created collection '{COLLECTION}'")
    client.close()

    progress = {} if reset else load_progress()
    next_id  = max((v["last_id"] for v in progress.values()), default=-1) + 1

    for file_path in tqdm(files, desc="Ingesting files", unit="file"):
        fhash = file_hash(file_path)
        if fhash in progress:
            tqdm.write(f"  Skipping (already indexed): {file_path.name}")
            continue

        tqdm.write(f"  Processing: {file_path.name}")
        course = derive_course(file_path, materials_dir)
        pages  = extract_file(file_path)
        if not pages:
            tqdm.write(f"    No extractable text — skipping")
            continue

        chunks = chunk_pages(pages, file_path.name, course)
        tqdm.write(f"    {len(pages)} pages → {len(chunks)} chunks")

        texts      = [c["text"] for c in chunks]
        embeddings = model.encode(texts, batch_size=64, show_progress_bar=False)

        points = [
            PointStruct(
                id=next_id + i,
                vector=emb.tolist(),
                payload=chunks[i],
            )
            for i, emb in enumerate(embeddings)
        ]

        # Fresh connection per file to avoid gRPC keepalive issues
        client = make_client()
        batch_sz = 64
        for b in range(0, len(points), batch_sz):
            client.points.upsert(COLLECTION, points[b:b + batch_sz])
        client.close()

        progress[fhash] = {
            "file":    file_path.name,
            "chunks":  len(chunks),
            "last_id": next_id + len(chunks) - 1,
        }
        next_id += len(chunks)
        save_progress(progress)

    client = make_client()
    total = client.points.count(COLLECTION)
    client.close()
    print(f"\nDone. {total} chunks indexed in '{COLLECTION}'.")
    print("Run `streamlit run app.py` to start the study assistant.")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Ingest documents into StudyVector")
    parser.add_argument("--dir",   default="materials", help="Folder with PDFs (default: ./materials)")
    parser.add_argument("--reset", action="store_true",  help="Wipe collection and re-ingest everything")
    args = parser.parse_args()

    materials_dir = Path(args.dir).expanduser().resolve()
    if not materials_dir.exists():
        print(f"Directory not found: {materials_dir}")
        sys.exit(1)

    ingest(materials_dir, reset=args.reset)
